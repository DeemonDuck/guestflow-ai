"""Generate the GuestFlow AI pitch deck as a .pptx.

    python presentation/generate_deck.py

Produces presentation/GuestFlow_AI.pptx. Edit the SLIDES data below to tweak
content, then re-run.
"""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Brand palette
NAVY = RGBColor(0x0F, 0x2A, 0x47)
TEAL = RGBColor(0x12, 0x9A, 0x8F)
LIGHT = RGBColor(0xF2, 0xF5, 0xF7)
GREY = RGBColor(0x44, 0x4A, 0x52)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "GuestFlow_AI.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _frame(box, text, size, color, bold=False, align=PP_ALIGN.LEFT, font="Segoe UI"):
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return tf


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def title_slide(title, subtitle, author):
    s = prs.slides.add_slide(BLANK)
    _bg(s, NAVY)
    bar = s.shapes.add_shape(1, Inches(0), Inches(2.7), SW, Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()
    _frame(s.shapes.add_textbox(Inches(0.8), Inches(2.9), Inches(11.7), Inches(1.4)),
           title, 54, WHITE, bold=True)
    _frame(s.shapes.add_textbox(Inches(0.85), Inches(4.2), Inches(11.7), Inches(0.8)),
           subtitle, 24, TEAL)
    _frame(s.shapes.add_textbox(Inches(0.85), Inches(6.4), Inches(11.7), Inches(0.6)),
           author, 16, LIGHT)


def bullet_slide(title, bullets, footer=None):
    s = prs.slides.add_slide(BLANK)
    _bg(s, WHITE)
    head = s.shapes.add_shape(1, Inches(0), Inches(0), SW, Inches(1.15))
    head.fill.solid(); head.fill.fore_color.rgb = NAVY; head.line.fill.background()
    _frame(s.shapes.add_textbox(Inches(0.6), Inches(0.28), Inches(12), Inches(0.7)),
           title, 30, WHITE, bold=True)
    accent = s.shapes.add_shape(1, Inches(0), Inches(1.15), SW, Inches(0.06))
    accent.fill.solid(); accent.fill.fore_color.rgb = TEAL; accent.line.fill.background()

    box = s.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(5.4))
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        lvl = 0
        text = b
        if isinstance(b, tuple):
            text, lvl = b
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        p.space_after = Pt(10)
        r = p.add_run()
        r.text = ("• " if lvl == 0 else "– ") + text
        r.font.size = Pt(20 if lvl == 0 else 17)
        r.font.color.rgb = NAVY if lvl == 0 else GREY
        r.font.name = "Segoe UI"
        r.font.bold = (lvl == 0)
    if footer:
        fb = s.shapes.add_textbox(Inches(0.7), Inches(6.85), Inches(12), Inches(0.5))
        _frame(fb, footer, 13, TEAL, bold=True)


def flow_slide(title, steps):
    s = prs.slides.add_slide(BLANK)
    _bg(s, WHITE)
    head = s.shapes.add_shape(1, Inches(0), Inches(0), SW, Inches(1.15))
    head.fill.solid(); head.fill.fore_color.rgb = NAVY; head.line.fill.background()
    _frame(s.shapes.add_textbox(Inches(0.6), Inches(0.28), Inches(12), Inches(0.7)),
           title, 30, WHITE, bold=True)
    accent = s.shapes.add_shape(1, Inches(0), Inches(1.15), SW, Inches(0.06))
    accent.fill.solid(); accent.fill.fore_color.rgb = TEAL; accent.line.fill.background()

    top = 2.4
    n = len(steps)
    box_w = 2.45
    gap = (13.333 - 0.6 * 2 - box_w * 0 - 0) / n
    x = 0.55
    for i, step in enumerate(steps):
        bx = x + i * gap + (gap - box_w) / 2
        shp = s.shapes.add_shape(5, Inches(bx), Inches(top), Inches(box_w), Inches(1.5))
        shp.fill.solid()
        shp.fill.fore_color.rgb = TEAL if i % 2 == 0 else NAVY
        shp.line.color.rgb = NAVY
        tf = shp.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = step
        r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE
        r.font.name = "Segoe UI"
        if i < n - 1:
            ar = s.shapes.add_textbox(Inches(bx + box_w - 0.05), Inches(top + 0.45), Inches(0.5), Inches(0.6))
            _frame(ar, "→", 26, NAVY, bold=True, align=PP_ALIGN.CENTER)
    cap = s.shapes.add_textbox(Inches(0.7), Inches(4.6), Inches(12), Inches(1.8))
    _frame(cap,
           "Most hotel AI stops at the first box. GuestFlow runs the whole chain — "
           "and confirms the guest is satisfied at the end.",
           20, GREY)


# ---- Build the deck ----
title_slide(
    "GuestFlow AI",
    "The AI Front Desk Assistant for Hotels",
    "Developed by Ridham Taneja",
)

bullet_slide("The Problem", [
    "Reception is interrupted all day by repetitive guest questions (Wi-Fi, checkout, towels, cabs).",
    "Guest issues get lost between reception, housekeeping and maintenance.",
    "Managers find out about problems too late — often as a public 1-star review.",
    "Independent hotels are short-staffed and can't scale attention.",
], footer="Repetitive work + slow issue resolution + reputation risk = lost revenue.")

bullet_slide("What GuestFlow Does", [
    "An AI front-desk assistant that handles repetitive guest requests instantly.",
    ("Communicate — answers questions and creates tickets automatically.", 1),
    ("Operate — routes, escalates and tracks every issue to resolution.", 1),
    ("Improve — protects ratings and surfaces problems before they grow.", 1),
], footer="Positioned as an operations system, not just a chatbot.")

flow_slide("The Differentiator: Closing the Loop", [
    "Guest reports issue", "AI creates ticket", "Staff notified",
    "Escalate if stale", "Resolved", "Guest confirmed",
])

bullet_slide("Key Features", [
    "Conversational guest support (RAG over hotel FAQ + local LLM).",
    "Ticket lifecycle: open → in-progress → resolved, with auto follow-up.",
    "Escalation timer — stale issues alert the manager automatically.",
    "Guest profiles & preferences — returning guests recognised and personalised.",
    "Review management — compliant feedback capture + negative-feedback recovery.",
    "Owner analytics + proactive insights + daily manager digest.",
])

bullet_slide("Review Management — Done Legally", [
    "Every guest is invited to leave a public review — never only the happy ones.",
    "Negative feedback privately alerts the manager for fast service recovery.",
    "The alert never blocks or diverts a guest's public review.",
    ("\"Review gating\" (soliciting only happy guests) violates Google & FTC rules — "
     "GuestFlow is built to avoid it by design.", 1),
], footer="Better ratings, the compliant way.")

bullet_slide("The Moat: Proactive Insights", [
    "Systemic issues — \"3 'Maintenance' tickets on floor 3 in 24h: likely one root cause.\"",
    "Repeat rooms — one room generating many tickets.",
    "Feedback themes — a topic recurring across negative reviews.",
    "Delivered live in the dashboard and as a manager morning digest.",
], footer="Most hotel AI answers your guests — GuestFlow tells you what's about to break.")

bullet_slide("Architecture", [
    "Event-driven webhook ingestion with semantic intent routing.",
    "Multi-agent orchestration: Pre-Stay, In-Stay, Post-Stay agents.",
    "RAG retrieval (ChromaDB) over the hotel FAQ + local LLM (Ollama).",
    "Persistent state in SQLite; durable guest memory in a vector store.",
    "FastAPI backend + Streamlit operations dashboard.",
])

bullet_slide("Security & Trust", [
    "API-key authentication (constant-time), fail-closed in production.",
    "Input validation, email header-injection protection, rate limiting.",
    "Prompt-injection hardening — untrusted guest text is fenced.",
    "Right-to-erasure endpoint for guest PII (DPDP/GDPR-friendly).",
    "48-test automated suite; runs behind HTTPS in deployment.",
])

bullet_slide("Tech Stack", [
    "Backend: FastAPI (Python).",
    "Dashboard: Streamlit.",
    "AI: local LLM via Ollama (swappable to hosted APIs).",
    "RAG: ChromaDB vector retrieval.",
    "Data: SQLite (WAL) — Postgres for multi-hotel scale.",
    "Testing: pytest (48 tests).",
])

bullet_slide("Why a Hotel Pays for This", [
    "Cuts repetitive front-desk work so staff focus on real guests.",
    "Faster issue resolution → fewer bad experiences → better reviews.",
    "Catches operational problems before they cost a star.",
    ("\"We help your staff respond faster, reduce repetitive work, and turn "
     "satisfied guests into positive reviews.\"", 1),
])

bullet_slide("Roadmap", [
    "WhatsApp / SMS guest channel (via a Business Solution Provider).",
    "Multi-tenant support — per-hotel configuration & data isolation.",
    "Cloud deployment + hosted LLM option.",
    "PMS integration and deeper semantic search.",
], footer="Built and tested today; deployment sequenced around a real pilot.")

title_slide(
    "Thank You",
    "GuestFlow AI — what's about to become a problem, before it does.",
    "Ridham Taneja",
)

prs.save(OUT)
print(f"Saved deck: {OUT} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
